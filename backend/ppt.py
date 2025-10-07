from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from PIL import Image
import io
import base64
from pathlib import Path

def process_ppt_file(file_path):
    """Process PPT file and extract slides with content and thumbnails"""
    try:
        prs = Presentation(file_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_data = extract_slide_data(slide, i)
            slides_data.append(slide_data)
        
        return {
            "total_slides": len(slides_data),
            "slides": slides_data,
            "filename": os.path.basename(file_path)
        }
        
    except Exception as e:
        raise Exception(f"Error processing PPT: {str(e)}")

def extract_slide_data(slide, slide_index):
    """Extract text and metadata from a single slide"""
    slide_text = []
    notes = ""
    
    # Extract slide notes
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame:
            notes = notes_slide.notes_text_frame.text
    
    # Extract text from shapes
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_text.append(shape.text)
        
        # Handle table content
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        slide_text.append(cell.text)
    
    # Generate a mock thumbnail (in production, you'd generate actual slide images)
    thumbnail = generate_mock_thumbnail(slide_index)
    
    return {
        "slide_number": slide_index + 1,
        "title": f"Slide {slide_index + 1}",
        "content": "\n".join(slide_text),
        "notes": notes,
        "thumbnail": thumbnail,
        "shapes_count": len(slide.shapes)
    }

def generate_mock_thumbnail(slide_index):
    """Generate a mock thumbnail for demonstration"""
    # In production, use python-pptx to generate actual slide images
    # For now, return a placeholder image URL
    colors = ['#4361ee', '#06d6a0', '#7209b7', '#f72585', '#4cc9f0']
    color = colors[slide_index % len(colors)]
    
    # Create a simple colored rectangle as base64 image
    from PIL import Image, ImageDraw
    img = Image.new('RGB', (300, 200), color)
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 290, 190], outline='white', width=2)
    draw.text((150, 100), f"Slide {slide_index + 1}", fill='white', anchor='mm')
    
    buffered = io.BytesIO()
    img.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    
    return f"data:image/png;base64,{img_str}"

def get_slide_content(file_path, slide_index):
    """Get detailed content for a specific slide"""
    prs = Presentation(file_path)
    
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError("Invalid slide index")
    
    slide = prs.slides[slide_index]
    return extract_slide_data(slide, slide_index)

def extract_text_from_ppt(file_path):
    """Extract all text from PPT for AI processing"""
    prs = Presentation(file_path)
    full_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        full_text.append(f"Slide {i+1}: " + " | ".join(slide_text))
    
    return "\n".join(full_text)

def save_uploaded_ppt(file_content: bytes, filename: str) -> str:
    """Save uploaded PPT file and return file path"""
    upload_dir = Path("static/uploads")
    upload_dir.mkdir(exist_ok=True)
    
    file_path = upload_dir / filename
    with open(file_path, "wb") as buffer:
        buffer.write(file_content)
    
    return str(file_path)